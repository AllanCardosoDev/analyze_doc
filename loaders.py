import os
from time import sleep
import streamlit as st
import pandas as pd
import numpy as np
import io
import json
import re
from langchain_community.document_loaders import (
    WebBaseLoader,
    CSVLoader,
    PyPDFLoader,
    TextLoader
)
from fake_useragent import UserAgent
from docx import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_transformers import DoctranTextTranslator
from langchain_core.documents import Document as LangchainDocument
from pptx import Presentation
from markdown import markdown
from bs4 import BeautifulSoup

try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

try:
    import spacy
    SPACY_AVAILABLE = True
    # Carregamos apenas se necessário no código
except ImportError:
    SPACY_AVAILABLE = False

def carrega_site(url):
    """Carrega texto de um site usando WebBaseLoader."""
    documento = ""
    for i in range(5):
        try:
            os.environ["USER_AGENT"] = UserAgent().random
            loader = WebBaseLoader(url, raise_for_status=True)
            lista_documentos = loader.load()
            documento = "\n\n".join([doc.page_content for doc in lista_documentos])
            break
        except Exception as e:
            print(f"Erro ao carregar o site {i+1}: {e}")
            sleep(3)
    if not documento:
        return "⚠️ Não foi possível carregar o site."
    return documento

def carrega_youtube(video_id, proxy=None):
    """
    Carrega legendas de vídeos do YouTube com suporte a proxy.
    
    Args:
        video_id: ID ou URL do vídeo do YouTube
        proxy: String de proxy no formato 'http://usuário:senha@host:porta'
    """
    try:
        # Extrai o video_id de uma URL completa, se for fornecida
        if "youtube.com" in video_id or "youtu.be" in video_id:
            if "youtube.com/watch?v=" in video_id:
                video_id = video_id.split("youtube.com/watch?v=")[1].split("&")[0]
            elif "youtu.be/" in video_id:
                video_id = video_id.split("youtu.be/")[1].split("?")[0]
        
        # Importa diretamente o youtube_transcript_api para mais controle
        from youtube_transcript_api import YouTubeTranscriptApi
        
        # Configuração de proxy, se fornecido
        proxies = None
        if proxy:
            proxies = {'http': proxy, 'https': proxy}
        
        # Tenta primeiro com o idioma português, depois com inglês
        languages = ['pt', 'pt-BR', 'en']
        
        # Obtém as legendas diretamente, com suporte a proxy
        if proxies:
            from youtube_transcript_api.formatters import TextFormatter
            from youtube_transcript_api import TranscriptListFetcher
            
            fetcher = TranscriptListFetcher(proxies=proxies)
            transcript_list = fetcher.fetch_transcript_list(video_id)
            transcript = transcript_list.find_transcript(languages)
            transcripts = transcript.fetch()
        else:
            # Método padrão sem proxy
            transcripts = YouTubeTranscriptApi.get_transcript(video_id, languages=languages)
        
        # Formata as legendas em um texto contínuo
        if not transcripts:
            return "⚠️ Não foi possível encontrar legendas para este vídeo."
        
        # Combina os segmentos de legendas em um texto formatado com timestamp
        texto_completo = ""
        for entry in transcripts:
            seconds = int(entry.get('start', 0))
            minutes = seconds // 60
            seconds = seconds % 60
            timestamp = f"[{minutes:02d}:{seconds:02d}]"
            texto_completo += f"{timestamp} {entry.get('text', '')} "
        
        return texto_completo
    except Exception as e:
        mensagem_erro = str(e)
        if "IP" in mensagem_erro and "block" in mensagem_erro:
            return """❌ O YouTube está bloqueando as requisições do seu IP. Isso pode acontecer por:
            
1. Muitas requisições foram feitas em um curto período
2. Você está usando um IP de um provedor de nuvem (AWS, Google Cloud, etc.)

Soluções:
- Configure um proxy na aba 'Configurações'
- Espere algumas horas e tente novamente
- Use uma VPN ou outra rede para acessar
            """
        return f"❌ Erro ao carregar YouTube: {e}"

def carrega_csv(caminho):
    """Carrega dados de arquivos CSV com suporte a detecção de codificação."""
    try:
        # Tenta primeiro com utf-8
        df = pd.read_csv(caminho, encoding='utf-8')
    except UnicodeDecodeError:
        # Se falhar, tenta com latin-1
        try:
            df = pd.read_csv(caminho, encoding='latin-1')
        except:
            # Última tentativa com detecção automática
            import chardet
            with open(caminho, 'rb') as f:
                resultado = chardet.detect(f.read())
            df = pd.read_csv(caminho, encoding=resultado['encoding'])
    
    # Converte dataframe para texto
    info_basicas = f"CSV com {len(df)} linhas e {len(df.columns)} colunas.\n\n"
    
    # Análise estatística básica
    estatisticas = "Estatísticas básicas:\n"
    for coluna in df.columns:
        if pd.api.types.is_numeric_dtype(df[coluna]):
            estatisticas += f"Coluna '{coluna}':\n"
            estatisticas += f"  - Média: {df[coluna].mean():.2f}\n"
            estatisticas += f"  - Mínimo: {df[coluna].min()}\n"
            estatisticas += f"  - Máximo: {df[coluna].max()}\n"
    
    # Converte para texto tabular para visualização no chat
    df_str = df.head(50).to_string()
    
    return info_basicas + estatisticas + "\n\nAmostra de dados:\n" + df_str

def carrega_pdf(caminho, ocr=False):
    """Carrega e extrai texto de um PDF com suporte opcional a OCR."""
    try:
        loader = PyPDFLoader(caminho)
        documentos = loader.load()
        texto = "\n\n".join([doc.page_content for doc in documentos])
        
        # Se o texto extraído for muito pequeno e OCR está habilitado
        if len(texto.strip()) < 100 and ocr and OCR_AVAILABLE:
            import fitz  # PyMuPDF
            texto_completo = ""
            
            # Abre o PDF e extrai as imagens para OCR
            doc = fitz.open(caminho)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                texto_pagina = pytesseract.image_to_string(img, lang='por+eng')
                texto_completo += f"\n--- Página {page_num + 1} ---\n{texto_pagina}\n"
            
            return texto_completo if texto_completo.strip() else "⚠️ Não foi possível extrair texto do PDF mesmo com OCR."
        
        return texto if texto.strip() else "⚠️ O PDF não contém texto extraível. Tente ativar a opção de OCR."
    except Exception as e:
        return f"❌ Erro ao carregar PDF: {e}"

def carrega_txt(caminho):
    """Carrega e extrai texto de um arquivo TXT."""
    try:
        with open(caminho, 'r', encoding='utf-8') as file:
            return file.read()
    except UnicodeDecodeError:
        # Se falhar com UTF-8, tenta com Latin-1
        try:
            with open(caminho, 'r', encoding='latin-1') as file:
                return file.read()
        except Exception as e:
            return f"❌ Erro ao carregar TXT: {e}"
    except Exception as e:
        return f"❌ Erro ao carregar TXT: {e}"

def carrega_docx(caminho):
    """Carrega e extrai texto de um arquivo DOCX."""
    try:
        doc = Document(caminho)
        texto_completo = []
        
        # Extrair texto de parágrafos
        for para in doc.paragraphs:
            if para.text.strip():
                # Verifica o estilo do parágrafo para preservar formatação
                if para.style.name.startswith('Heading'):
                    texto_completo.append(f"## {para.text}")
                else:
                    texto_completo.append(para.text)
        
        # Extrair texto de tabelas
        for table in doc.tables:
            tabela_texto = []
            for i, row in enumerate(table.rows):
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    if i == 0:  # Assume primeira linha como cabeçalho
                        tabela_texto.append("| " + " | ".join(row_text) + " |")
                        tabela_texto.append("| " + " | ".join(["---"] * len(row_text)) + " |")
                    else:
                        tabela_texto.append("| " + " | ".join(row_text) + " |")
            
            if tabela_texto:
                texto_completo.append("\n".join(tabela_texto))
        
        resultado = "\n\n".join(texto_completo)
        
        # Adicionar metadados
        metadados = "# Metadados do Documento\n"
        try:
            props = doc.core_properties
            metadados += f"- Título: {props.title or 'Não disponível'}\n"
            metadados += f"- Autor: {props.author or 'Não disponível'}\n"
            metadados += f"- Criado em: {props.created or 'Não disponível'}\n"
            metadados += f"- Modificado em: {props.modified or 'Não disponível'}\n"
            resultado = metadados + "\n\n" + resultado
        except:
            pass
        
        return resultado if resultado.strip() else "⚠️ O documento DOCX não contém texto extraível."
    except Exception as e:
        return f"❌ Erro ao carregar DOCX: {e}"

def carrega_xlsx(caminho):
    """Carrega e extrai dados de um arquivo Excel."""
    try:
        # Lê todas as planilhas do arquivo
        xls = pd.ExcelFile(caminho)
        sheet_names = xls.sheet_names
        
        texto_completo = f"# Arquivo Excel com {len(sheet_names)} planilhas\n\n"
        
        # Processa cada planilha
        for sheet_name in sheet_names:
            df = pd.read_excel(caminho, sheet_name=sheet_name)
            
            texto_completo += f"## Planilha: {sheet_name}\n"
            texto_completo += f"- Linhas: {len(df)}\n"
            texto_completo += f"- Colunas: {len(df.columns)}\n\n"
            
            # Estatísticas básicas para colunas numéricas
            colunas_numericas = df.select_dtypes(include=[np.number]).columns
            if len(colunas_numericas) > 0:
                texto_completo += "### Estatísticas básicas\n"
                for coluna in colunas_numericas[:5]:  # Limitamos a 5 colunas para não sobrecarregar
                    texto_completo += f"**{coluna}**: "
                    texto_completo += f"Média: {df[coluna].mean():.2f}, "
                    texto_completo += f"Mín: {df[coluna].min()}, "
                    texto_completo += f"Máx: {df[coluna].max()}\n"
                texto_completo += "\n"
            
            # Amostra dos dados (primeiras 10 linhas)
            texto_completo += "### Amostra de dados\n"
            texto_completo += df.head(10).to_markdown() + "\n\n"
            
        return texto_completo
    except Exception as e:
        return f"❌ Erro ao carregar arquivo Excel: {e}"

def carrega_pptx(caminho):
    """Carrega e extrai texto de arquivos PowerPoint."""
    try:
        prs = Presentation(caminho)
        
        # Extrai texto de todos os slides
        texto_completo = f"# Apresentação PowerPoint com {len(prs.slides)} slides\n\n"
        
        for i, slide in enumerate(prs.slides):
            texto_completo += f"## Slide {i+1}\n"
            
            # Tenta extrair o título do slide
            if slide.shapes.title and slide.shapes.title.text:
                texto_completo += f"### {slide.shapes.title.text}\n"
            
            # Extrai texto de todas as formas no slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    if shape.text != slide.shapes.title.text if slide.shapes.title else True:
                        texto_completo += f"- {shape.text.strip()}\n"
            
            texto_completo += "\n"
        
        return texto_completo
    except Exception as e:
        return f"❌ Erro ao carregar arquivo PowerPoint: {e}"

def carrega_json(caminho):
    """Carrega e processa um arquivo JSON."""
    try:
        with open(caminho, 'r', encoding='utf-8') as file:
            data = json.load(file)
        
        # Converte para uma string formatada
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception as e:
        return f"❌ Erro ao carregar arquivo JSON: {e}"

def carrega_markdown(caminho):
    """Carrega e processa um arquivo Markdown."""
    try:
        with open(caminho, 'r', encoding='utf-8') as file:
            md_text = file.read()
        
        # Retornamos o texto markdown diretamente
        return md_text
    except Exception as e:
        return f"❌ Erro ao carregar arquivo Markdown: {e}"

def carrega_html(caminho):
    """Carrega e extrai texto de um arquivo HTML."""
    try:
        with open(caminho, 'r', encoding='utf-8') as file:
            html_content = file.read()
        
        # Usa BeautifulSoup para extrair o texto
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove scripts e estilos
        for script in soup(["script", "style"]):
            script.extract()
        
        # Extrai texto
        text = soup.get_text()
        
        # Limpa espaço em branco
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
        
        return text
    except Exception as e:
        return f"❌ Erro ao carregar arquivo HTML: {e}"

def carrega_imagem(caminho, nivel_ocr="simples"):
    """Extrai texto de uma imagem usando OCR."""
    if not OCR_AVAILABLE:
        return "⚠️ OCR não disponível. Instale as bibliotecas pytesseract e pillow."
    
    try:
        img = Image.open(caminho)
        
        # Configurações de OCR baseadas no nível selecionado
        ocr_config = ''
        if nivel_ocr == "avançado":
            ocr_config = '--oem 1 --psm 6'
        elif nivel_ocr == "completo":
            ocr_config = '--oem 1 --psm 6 -c preserve_interword_spaces=1'
        
        # Extrai o texto
        texto = pytesseract.image_to_string(img, lang='por+eng', config=ocr_config)
        
        # Se não houver texto suficiente, tenta pré-processamento
        if len(texto.strip()) < 100 and nivel_ocr in ["avançado", "completo"]:
            from PIL import ImageEnhance, ImageFilter
            
            # Pré-processamento para melhorar OCR
            img = img.convert('L')  # Converte para escala de cinza
            img = img.filter(ImageFilter.SHARPEN)  # Aumenta nitidez
            enhancer = ImageEnhance.Contrast(img)
            img = enhancer.enhance(2.0)  # Aumenta contraste
            
            texto = pytesseract.image_to_string(img, lang='por+eng', config=ocr_config)
        
        if not texto.strip():
            return "⚠️ Não foi possível extrair texto da imagem."
        
        return texto
    except Exception as e:
        return f"❌ Erro ao processar imagem: {e}"

def gera_resumo(texto, max_length=1000):
    """
    Gera um resumo automático do texto usando chunking e extração de frases importantes.
    Esta é uma implementação básica. Para resultados melhores, use os modelos de LLM.
    """
    try:
        # Verifica se o texto já é curto o suficiente
        if len(texto) <= max_length:
            return texto
            
        # Converter para documento Langchain
        doc = LangchainDocument(page_content=texto)
        
        # Dividir em chunks para processamento
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        chunks = text_splitter.split_documents([doc])
        
        # Extrair frases importantes usando uma abordagem baseada em heurísticas
        # Esta é uma implementação simplificada sem ML
        frases_importantes = []
        palavras_chave = extrair_palavras_chave(texto, 20)
        
        for chunk in chunks:
            # Dividir em frases
            sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', chunk.page_content) if s.strip()]
            
            # Pontuar frases com base em heurísticas
            scored_sentences = []
            for sentence in sentences:
                # Pontuação baseada no comprimento (não queremos frases muito curtas nem muito longas)
                length_score = min(len(sentence.split()) / 20.0, 1.0)
                
                # Pontuação baseada em palavras-chave
                keyword_score = sum(1 for word in palavras_chave if word.lower() in sentence.lower()) / len(palavras_chave)
                
                # Posição na sequência (frases iniciais são mais importantes)
                position_score = 1.0 if sentences.index(sentence) < 2 else 0.5
                
                # Score final ponderado
                final_score = (0.4 * length_score) + (0.4 * keyword_score) + (0.2 * position_score)
                
                scored_sentences.append((sentence, final_score))
            
            # Ordenar por pontuação
            scored_sentences.sort(key=lambda x: x[1], reverse=True)
            
            # Pegar as 3 frases mais bem pontuadas
            top_sentences = [s[0] for s in scored_sentences[:3]]
            frases_importantes.extend(top_sentences)
        
        # Limitar o número total de frases baseado no comprimento desejado
        resumo = ". ".join(frases_importantes)
        
        # Se ainda estiver muito longo, faz um corte simples
        if len(resumo) > max_length:
            # Tenta cortar em um ponto final para manter a coerência
            ultimo_ponto = resumo[:max_length].rfind('.')
            if ultimo_ponto > max_length * 0.7:  # Se estiver nos últimos 30%
                resumo = resumo[:ultimo_ponto + 1]
            else:
                resumo = resumo[:max_length] + "..."
        
        return resumo
    except Exception as e:
        print(f"Erro ao gerar resumo: {e}")
        return texto[:max_length] + "..."  # Fallback para um corte simples

def traduz_texto(texto, idioma_destino='pt'):
    """
    Traduz o texto para o idioma especificado.
    Requer que um modelo de LLM esteja disponível.
    """
    try:
        # Criar um documento Langchain
        doc = LangchainDocument(page_content=texto)
        
        # Configurar o tradutor
        translator = DoctranTextTranslator(
            language=idioma_destino
        )
        
        # Traduzir o documento
        documentos_traduzidos = translator.transform_documents([doc])
        
        return documentos_traduzidos[0].page_content
    except Exception as e:
        print(f"Erro ao traduzir texto: {e}")
        return f"Não foi possível traduzir o texto: {e}"

def extrai_entidades(texto):
    """Extrai entidades (pessoas, organizações, locais, etc.) do texto."""
    if not SPACY_AVAILABLE:
        return ["Biblioteca spaCy não disponível para extração de entidades"]
    
    try:
        # Carrega o modelo spaCy para português
        nlp = spacy.load("pt_core_news_sm")
        
        # Se o texto for muito longo, dividimos para processamento
        MAX_LENGTH = 100000  # spaCy tem limite de tamanho
        if len(texto) > MAX_LENGTH:
            trechos = [texto[i:i+MAX_LENGTH] for i in range(0, len(texto), MAX_LENGTH)]
        else:
            trechos = [texto]
            
        entidades_dict = {
            "PER": [],  # Pessoas
            "ORG": [],  # Organizações
            "LOC": [],  # Locais
            "GPE": [],  # Países, cidades
            "MISC": []  # Outros
        }
        
        # Processa cada trecho
        for trecho in trechos:
            doc = nlp(trecho)
            for ent in doc.ents:
                if ent.label_ in ["PERSON", "PER"]:
                    entidades_dict["PER"].append(ent.text)
                elif ent.label_ in ["ORG", "ORGANIZATION"]:
                    entidades_dict["ORG"].append(ent.text)
                elif ent.label_ in ["LOC", "LOCATION"]:
                    entidades_dict["LOC"].append(ent.text)
                elif ent.label_ in ["GPE", "COUNTRY", "CITY"]:
                    entidades_dict["GPE"].append(ent.text)
                else:
                    entidades_dict["MISC"].append(ent.text)
        
        # Remove duplicatas e limita o número de entidades
        for key in entidades_dict:
            entidades_dict[key] = list(set(entidades_dict[key]))[:20]  # Limita a 20 entidades por categoria
        
        # Retorna as entidades como lista plana
        entidades_flat = []
        for key in ["PER", "ORG", "LOC", "GPE"]:
            entidades_flat.extend(entidades_dict[key])
        
        return entidades_flat
    except Exception as e:
        print(f"Erro ao extrair entidades: {e}")
        return []

def analisa_sentimento(texto):
    """Analisa o sentimento (positivo, negativo, neutro) do texto."""
    try:
        from textblob import TextBlob
        
        # Divide o texto em parágrafos para análise
        paragrafos = texto.split('\n\n')
        
        # Analisa cada parágrafo
        polaridades = []
        for paragrafo in paragrafos:
            if len(paragrafo.strip()) > 20:  # Ignora parágrafos muito curtos
                blob = TextBlob(paragrafo)
                polaridades.append(blob.sentiment.polarity)
        
        # Calcula a polaridade média
        if polaridades:
            media_polaridade = sum(polaridades) / len(polaridades)
            
            # Classifica o sentimento
            if media_polaridade > 0.2:
                return "Positivo"
            elif media_polaridade < -0.2:
                return "Negativo"
            else:
                return "Neutro"
        else:
            return "Neutro"
    except Exception as e:
        print(f"Erro ao analisar sentimento: {e}")
        return "Não foi possível analisar o sentimento"

def detecta_topicos(texto):
    """Detecta os principais tópicos do texto usando LDA ou métodos similares."""
    try:
        from sklearn.feature_extraction.text import CountVectorizer
        from sklearn.decomposition import LatentDirichletAllocation
        import nltk
        from nltk.corpus import stopwords
        
        # Certifica-se que temos os recursos necessários
        try:
            nltk.data.find('corpora/stopwords')
        except LookupError:
            nltk.download('stopwords')
        
        # Stopwords em português
        stop_words = set(stopwords.words('portuguese'))
        
        # Divide o texto em parágrafos
        paragrafos = [p for p in texto.split('\n\n') if len(p.strip()) > 50]
        
        if len(paragrafos) < 3:
            # Para documentos curtos, usamos uma abordagem baseada em frequência
            palavras = re.findall(r'\b\w{3,15}\b', texto.lower())
            palavras = [p for p in palavras if p not in stop_words and not p.isdigit()]
            
            # Conta a frequência
            frequencia = {}
            for palavra in palavras:
                frequencia[palavra] = frequencia.get(palavra, 0) + 1
            
            # Ordena por frequência
            palavras_ordenadas = sorted(frequencia.items(), key=lambda x: x[1], reverse=True)
            
            # Retorna os 10 primeiros como tópicos
            return [palavra for palavra, _ in palavras_ordenadas[:10]]
        else:
            # Para documentos mais longos, usamos LDA
            vectorizer = CountVectorizer(max_df=0.95, min_df=2, stop_words=list(stop_words))
            dtm = vectorizer.fit_transform(paragrafos)
            
            # Número de tópicos
            n_topicos = min(5, len(paragrafos) // 2)
            
            lda = LatentDirichletAllocation(n_components=n_topicos, random_state=42)
            lda.fit(dtm)
            
            # Extrai os termos mais importantes de cada tópico
            feature_names = vectorizer.get_feature_names_out()
            topicos = []
            
            for topic_idx, topic in enumerate(lda.components_):
                topico = " ".join([feature_names[i] for i in topic.argsort()[:-6:-1]])
                topicos.append(topico)
            
            return topicos
    except Exception as e:
        print(f"Erro ao detectar tópicos: {e}")
        return ["Não foi possível detectar tópicos"]

def classifica_documento(texto):
    """Classifica o tipo/categoria do documento."""
    # Implementação simplificada baseada em regras
    texto_lower = texto.lower()
    
    # Palavras-chave para categorias comuns
    categorias = {
        "Contrato": ["contrato", "acordo", "partes", "cláusula", "obrigações", "rescisão"],
        "Artigo Científico": ["abstract", "metodologia", "resultados", "conclusão", "referências", "estudo"],
        "Notícia": ["notícia", "jornal", "reportagem", "ontem", "hoje", "declarou"],
        "Manual Técnico": ["manual", "instrução", "passos", "procedimento", "técnico", "equipamento"],
        "Email": ["prezado", "atenciosamente", "cordialmente", "e-mail", "att", "em anexo"],
        "Relatório": ["relatório", "análise", "avaliação", "período", "desempenho", "constatou-se"],
        "Apresentação": ["slide", "apresentação", "tópico", "bullet", "próximo slide"],
        "Legislação": ["lei", "artigo", "parágrafo", "inciso", "decreto", "legislação"]
    }
    
    # Conta ocorrências de palavras-chave
    pontuacao = {cat: 0 for cat in categorias}
    for categoria, palavras in categorias.items():
        for palavra in palavras:
            pontuacao[categoria] += texto_lower.count(palavra)
    
    # Verifica formatos específicos
    if "```json" in texto or texto.strip().startswith("{") and texto.strip().endswith("}"):
        return "Arquivo JSON"
    elif "```html" in texto or ("<html" in texto_lower and "</html>" in texto_lower):
        return "Documento HTML"
    elif "```sql" in texto or "select " in texto_lower and " from " in texto_lower:
        return "Script SQL"
    elif "```python" in texto or "def " in texto and ":" in texto:
        return "Código Python"
    
    # Retorna a categoria com maior pontuação
    if any(pontuacao.values()):
        melhor_categoria = max(pontuacao.items(), key=lambda x: x[1])
        if melhor_categoria[1] > 0:
            return melhor_categoria[0]
    
    # Se não conseguir classificar
    return "Documento Texto Geral"

def extrai_tabelas(texto):
    """Extrai tabelas do texto."""
    # Busca por padrões de tabela em markdown ou texto
    tabelas = []
    
    # Procura tabelas markdown
    padrao_md = r'\|[^\n]+\|\n\|[\s\-:\|]+\|\n(?:\|[^\n]+\|\n)+'
    tabelas_md = re.findall(padrao_md, texto)
    
    for tabela in tabelas_md:
        tabelas.append({"tipo": "markdown", "conteudo": tabela})
    
    # Procura tabelas em formato ASCII/texto
    linhas = texto.split('\n')
    i = 0
    while i < len(linhas):
        # Verifica se a linha tem potencial para ser cabeçalho de tabela
        if '+' in linhas[i] and '-' in linhas[i]:
            inicio = i
            i += 1
            tabela_ascii = [linhas[inicio]]
            
            # Coleta linhas que parecem fazer parte da tabela
            while i < len(linhas) and ('+' in linhas[i] or '|' in linhas[i]):
                tabela_ascii.append(linhas[i])
                i += 1
            
            # Se encontramos pelo menos 3 linhas, pode ser uma tabela
            if len(tabela_ascii) >= 3:
                tabelas.append({
                    "tipo": "ascii",
                    "conteudo": "\n".join(tabela_ascii)
                })
        else:
            i += 1
    
    return tabelas

def extrai_dados_financeiros(texto):
    """Extrai valores monetários, porcentagens e outros dados financeiros."""
    dados_financeiros = {
        "valores_monetarios": [],
        "porcentagens": [],
        "datas": []
    }
    
    # Extrai valores monetários (R$, $, €)
    padrao_moeda = r'(?:R\$|\$|€)\s*\d+(?:[.,]\d+)*(?:\s*(?:mil|milhão|milhões|bilhão|bilhões))?'
    valores_monetarios = re.findall(padrao_moeda, texto)
    dados_financeiros["valores_monetarios"] = valores_monetarios[:20]  # Limita a 20 valores
    
    # Extrai porcentagens
    padrao_porcentagem = r'\d+(?:[.,]\d+)*\s*%'
    porcentagens = re.findall(padrao_porcentagem, texto)
    dados_financeiros["porcentagens"] = porcentagens[:20]
    
    # Extrai datas em vários formatos
    padrao_data = r'(?:\d{1,2}/\d{1,2}/\d{2,4}|\d{1,2} de (?:janeiro|fevereiro|março|abril|maio|junho|julho|agosto|setembro|outubro|novembro|dezembro) de \d{2,4})'
    datas = re.findall(padrao_data, texto)
    dados_financeiros["datas"] = datas[:20]
    
    return dados_financeiros

def extrair_palavras_chave(texto, num_palavras=10):
    """
    Extrai as palavras-chave mais relevantes de um texto.
    Uma implementação simples baseada em frequência.
    """
    try:
        # Stopwords em português
        try:
            import nltk
            from nltk.corpus import stopwords
            try:
                nltk.data.find('corpora/stopwords')
            except LookupError:
                nltk.download('stopwords')
            stop_words = set(stopwords.words('portuguese'))
        except:
            # Fallback se NLTK não estiver disponível
            stop_words = set(['a', 'ao', 'aos', 'aquela', 'aquelas', 'aquele', 'aqueles', 'aquilo', 'as', 'até', 'com', 'como', 'da', 'das', 'de', 'dela', 'delas', 'dele', 'deles', 'depois', 'do', 'dos', 'e', 'ela', 'elas', 'ele', 'eles', 'em', 'entre', 'era', 'eram', 'éramos', 'essa', 'essas', 'esse', 'esses', 'esta', 'estas', 'este', 'estes', 'eu', 'foi', 'fomos', 'for', 'foram', 'forem', 'fui', 'há', 'isso', 'isto', 'já', 'lhe', 'lhes', 'mais', 'mas', 'me', 'mesmo', 'meu', 'meus', 'minha', 'minhas', 'muito', 'na', 'não', 'nas', 'nem', 'no', 'nos', 'nós', 'nossa', 'nossas', 'nosso', 'nossos', 'num', 'numa', 'o', 'os', 'ou', 'para', 'pela', 'pelas', 'pelo', 'pelos', 'por', 'qual', 'quando', 'que', 'quem', 'são', 'se', 'seja', 'sejam', 'sejamos', 'sem', 'será', 'serão', 'serei', 'seremos', 'seria', 'seriam', 'seríamos', 'seu', 'seus', 'só', 'somos', 'sou', 'sua', 'suas', 'também', 'te', 'tem', 'tém', 'temos', 'tenho', 'teu', 'teus', 'tu', 'tua', 'tuas', 'um', 'uma', 'você', 'vocês', 'vos'])
            
        # Encontra todas as palavras no texto
        palavras = re.findall(r'\b[a-zA-ZáàâãéèêíïóôõöúçñÁÀÂÃÉÈÊÍÏÓÔÕÖÚÇÑ]{3,15}\b', texto.lower())
        
        # Remove stopwords e conta frequência
        palavras_filtradas = [palavra for palavra in palavras if palavra not in stop_words]
        
        # Conta a frequência
        frequencia = {}
        for palavra in palavras_filtradas:
            frequencia[palavra] = frequencia.get(palavra, 0) + 1
        
        # Ordena por frequência
        palavras_ordenadas = sorted(frequencia.items(), key=lambda x: x[1], reverse=True)
        
        # Retorna as X palavras mais frequentes
        return [palavra for palavra, _ in palavras_ordenadas[:num_palavras]]
    except Exception as e:
        print(f"Erro ao extrair palavras-chave: {e}")
        return []

def anonimizar_texto(texto, nivel="básico"):
    """
    Anonimiza informações pessoais no texto.
    Níveis: básico, intermediário, completo
    """
    texto_anonimizado = texto
    
    # Anonimiza CPFs
    texto_anonimizado = re.sub(r'\d{3}\.\d{3}\.\d{3}-\d{2}', '***.***.***-**', texto_anonimizado)
    
    # Anonimiza e-mails
    texto_anonimizado = re.sub(r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}', '****@****.com', texto_anonimizado)
    
    # Anonimiza telefones em vários formatos
    texto_anonimizado = re.sub(r'\(\d{2}\)\s*\d{4,5}-\d{4}', '(**)****-****', texto_anonimizado)
    texto_anonimizado = re.sub(r'\d{4,5}-\d{4}', '****-****', texto_anonimizado)
    
    if nivel in ["intermediário", "completo"]:
        # Anonimização intermediária inclui endereços e RG
        texto_anonimizado = re.sub(r'\d{1,5}\.\d{3}\.\d{3}-\d{1}', '***.***-*', texto_anonimizado)  # RG
        
        # Tentativa de anonimizar endereços (simplificada)
        termos_endereco = ["rua ", "avenida ", "av ", "alameda ", "praça ", "travessa ", "rodovia "]
        for termo in termos_endereco:
            padrao = f"{termo}[a-zA-Z\s]+,?\s*\d+"
            texto_anonimizado = re.sub(padrao, f"{termo}****", texto_anonimizado, flags=re.IGNORECASE)
    
    if nivel == "completo":
        # Anonimização completa inclui nomes próprios
        # Esta é uma abordagem muito simplificada e pode gerar falsos positivos
        # Uma implementação real precisaria de NER (Named Entity Recognition)
        
        if SPACY_AVAILABLE:
            try:
                nlp = spacy.load("pt_core_news_sm")
                doc = nlp(texto)
                
                # Cria um dicionário de substituições
                substituicoes = {}
                
                for ent in doc.ents:
                    if ent.label_ in ["PERSON", "PER"]:
                        # Substitui nomes mantendo o primeiro caractere
                        substituicoes[ent.text] = f"{ent.text[0]}****"
                
                # Aplica as substituições
                for original, substituicao in substituicoes.items():
                    texto_anonimizado = texto_anonimizado.replace(original, substituicao)
            except Exception as e:
                print(f"Erro ao anonimizar nomes: {e}")
    
    return texto_anonimizado
